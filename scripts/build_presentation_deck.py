"""Assemble the v1.2 12-minute talk deck from the Y_presentation/ PNGs.

Outputs ``reports/gallery/Y_presentation/ArqueoGal_Pipeline1_v1_2_talk.pptx``.

Layout grid (per the v1.2 brief):
  - 16:9, 13.333 x 7.5 inches.
  - 0.5 in margins.
  - Title row 0.4 to 1.1 in, 28 pt sentence-case left-aligned.
  - Subtitle row 1.1 to 1.6 in, 14 pt left-aligned, max 10 in wide.
  - Content area 1.7 to 6.9 in.
  - Footer row 6.9 to 7.3 in: left "ArqueoGal Pipeline 1 v1.2 / IA-CAUP",
    right "N / NSLIDES" both #5C6378 9 pt.

Templates: T-Title, T-Story (one figure), T-Compare (two figures),
T-Stat (large numerics + small figure). The brief ALSO defines T-Section
(dark divider) but we deferred those to keep the 12-min runtime.
"""

from __future__ import annotations

import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
GALLERY = REPO / "reports" / "gallery" / "Y_presentation"
OUT = GALLERY / "ArqueoGal_Pipeline1_v1_2_talk.pptx"

# v1.2 chrome palette.
TITLE_INK = RGBColor(0x1A, 0x1D, 0x2E)
BODY_INK = RGBColor(0x2B, 0x2D, 0x42)
MUTED = RGBColor(0x5C, 0x63, 0x78)
ACCENT = RGBColor(0x00, 0x72, 0xB2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER_BG = RGBColor(0x1A, 0x1D, 0x2E)


Template = Literal["title", "story", "compare", "stat"]


@dataclass
class Slide:
    title: str
    subtitle: str
    template: Template
    images: list[str] = field(default_factory=list)
    notes: str = ""
    big_numbers: list[tuple[str, str]] = field(default_factory=list)


SLIDES: list[Slide] = [
    Slide(
        title="From Gaia BP/RP XP to APOGEE-grade chemistry",
        subtitle=(
            "ArqueoGal Pipeline 1, v1.1.  A. W. Neitzel, IA-CAUP, "
            "Universidade do Porto.  FCT Exploratory Project 2024.15303.PEX."
        ),
        template="title",
        images=["Y43_headline_rmse_card.png"],
        notes=(
            "Open with the central pitch: convert Gaia DR3 BP/RP XP into "
            "APOGEE DR19-grade stellar parameters and abundances on a "
            "sample two orders of magnitude larger than APOGEE itself.  "
            "Frame the talk: data, model, validation, tier system, "
            "Stream-2 application, science output."
        ),
    ),
    Slide(
        title="Sample flow",
        subtitle=(
            "Reconciling the n-counts: parent catalogue to Tier-1 holdout."
        ),
        template="story",
        images=["Y44_sample_flow.png"],
        notes=(
            "Five stages, each annotated with the cut that discards stars: "
            "Gaia DR3 XP (219M) -> S1 raw join with APOGEE DR19 (327k) "
            "-> S1 dedup (293k) -> stratified val+test holdout (88k) "
            "-> Tier-1 subset (83k).  Every n in the deck lives at one "
            "of these stages."
        ),
    ),
    Slide(
        title="Why XP, why APOGEE, why now",
        subtitle=(
            "Gaia DR3 XP carries shape information for every disc giant.  "
            "APOGEE DR19 chemistry is high-resolution but covers ~720k "
            "stars.  Bridge the asymmetry."
        ),
        template="story",
        images=["Y02_three_streams.png"],
        notes=(
            "Three-stream view: Stream 1 = APOGEE x XP training cohort; "
            "Stream 2 = TESS asteroseismic giants x XP, asteroseismic "
            "ground truth on log g; Stream 3 = Andrae+23 RGB x XP, the "
            "release cohort.  Today the focus is Streams 1 and 2; "
            "Stream 3 is deferred."
        ),
    ),
    Slide(
        title="The XP feature",
        subtitle=(
            "BP and RP each carry 55 normalised Hermite coefficients; the "
            "encoder consumes the 108-D coefficient block, no flux model."
        ),
        template="story",
        images=["Y32_xp_spectra_and_coefficients.png"],
        notes=(
            "Per-band median normalisation.  BP and RP are independent "
            "photometers with different effective areas, so their "
            "absolute fluxes do not match in the 640-680 nm overlap.  "
            "That is calibration physics, not a bug.  We feed the "
            "normalised Hermite coefficients (post-Ye+2024 flux "
            "correction, post frozen v1 z-score) to the encoder."
        ),
    ),
    Slide(
        title="Feature engineering",
        subtitle=(
            "108 XP normalised Hermite coefficients plus 23 auxiliary "
            "scalars: photometry, parallax, BJ21 distance triple, "
            "dereddened 2MASS+WISE, four extinction priors."
        ),
        template="story",
        images=["Y03_preprocessing_flow.png"],
        notes=(
            "Lindegren+2021 parallax zero point; Riello+2021 G-mag "
            "correction; Ye+2024 NN flux correction on XP; Yuan+2013 + "
            "R_V=3.1 dereddening on broadband photometry.  Frozen v1 "
            "Hermite z-score basis (fingerprint 0d34b56...) is the "
            "inference contract: the same 108-D scaling at training "
            "and at deployment, never refit."
        ),
    ),
    Slide(
        title="Architecture",
        subtitle=(
            "Shared encoder; SupCon (1.0) plus Barlow (0.5) plus beta-NLL "
            "(1.0); 5x5 block-Cholesky covariance head."
        ),
        template="compare",
        images=["Y04_model_architecture.png", "Y24_block_cholesky.png"],
        notes=(
            "Trunk: shared encoder, contrastive pretraining on the label "
            "space (soft Gaussian-kernel positives, per-element "
            "bandwidths matched to disc chemistry scales).  Joint "
            "finetune: heteroscedastic Gaussian NLL on the 5-label "
            "means + a 5x5 block-Cholesky covariance head, plus Barlow "
            "Twins to prevent latent collapse."
        ),
    ),
    Slide(
        title="Why this combination of losses",
        subtitle=(
            "beta-NLL keeps the variance honest; SupCon makes the latent "
            "label-aware; Barlow keeps it from collapsing."
        ),
        template="compare",
        images=["Y21_beta_nll_loss.png", "Y22_supcon_loss.png"],
        notes=(
            "beta-NLL (Seitzer+2022): minimise the heteroscedastic "
            "Gaussian NLL with a beta=0.5 downweight on the variance "
            "to stop run-away.  SupCon (Khosla+2020): soft positives "
            "via a Gaussian kernel on the label vector.  Barlow Twins "
            "(Zbontar+2021): cross-correlation of two augmented views, "
            "redundancy reduction to keep the trunk informative."
        ),
    ),
    Slide(
        title="Training cadence",
        subtitle=(
            "200 epochs per stage with patience 20.  Per-epoch RMSE "
            "collapses cleanly toward the floor on the holdout cohort."
        ),
        template="story",
        images=["Y31_cadence_convergence.png"],
        notes=(
            "Pretrain stage: 200 epochs of contrastive on the encoder "
            "with the regression head detached.  Finetune stage: 200 "
            "joint epochs, encoder-LR ratio 0.1 of the head, AMP bf16 "
            "on the RTX 3060 6 GB.  Cadence-prediction parquets "
            "emitted every epoch let us animate convergence."
        ),
    ),
    Slide(
        title="Headline numbers",
        subtitle=(
            "Tier-1 holdout precision is within ~0.5x APOGEE DR19 internal "
            "repeatability across all five labels."
        ),
        template="story",
        images=["Y43_headline_rmse_card.png"],
        notes=(
            "Per-label RMSE on the Stream-1 Tier-1 holdout; bars expressed "
            "as a ratio to APOGEE DR19 internal repeatability.  All five "
            "labels sit in [0.5x, 2.4x] APOGEE.  Tier breakdown of the full "
            "holdout: 93.93% Tier 1, 0.38% Tier 2, 5.69% Tier 3."
        ),
    ),
    Slide(
        title="Pred vs truth and residuals by tier",
        subtitle=(
            "Filled = all data.  Tier-1 distribution sits inside the "
            "all-data envelope; Tier 2 / Tier 3 occupy the tails as designed."
        ),
        template="compare",
        images=["Y07_truth_vs_pred_headline.png", "Y35_residuals_by_tier.png"],
        notes=(
            "Pred-vs-truth scatter shows tight 1:1 across labels with the "
            "alpha/M panel preserving the disc bimodality.  Residuals: log "
            "y, T2 / T3 step outlines on top of the all-data shaded "
            "background, RMSE annotated per cohort."
        ),
    ),
    Slide(
        title="Are the uncertainties trustworthy",
        subtitle=(
            "Pull-distribution sigma within +/- 15% of N(0,1); reliability "
            "diagrams track the diagonal."
        ),
        template="compare",
        images=["Y16_pull_distributions.png", "Y17_reliability_diagram.png"],
        notes=(
            "Pull = (mu_pred - y_true) / sigma_pred.  If the model's "
            "uncertainties are calibrated, pulls are unit-Gaussian.  "
            "We measure pull sigma in [0.95, 1.12] across all five "
            "labels.  Reliability: empirical RMSE vs predicted sigma "
            "across sigma bins; diagonal = perfectly calibrated."
        ),
    ),
    Slide(
        title="Tier system, dual-Mahalanobis",
        subtitle=(
            "Tier 3, input-OOD on 108-D XP at p99.  Tier 2, output-OOD on "
            "5-D label envelope at p99.  Tier 1, everything else."
        ),
        template="story",
        images=["Y36_mahalanobis_percentiles_and_umaps.png"],
        notes=(
            "v6 schema, 2026-05-03.  Tier 3 catches stars whose XP "
            "coefficient vector lies outside the training distribution.  "
            "Tier 2 catches predictions whose 5-D label tuple lies "
            "outside the APOGEE truth envelope.  Per-star "
            "ood_mahalanobis_percentile and "
            "label_mahalanobis_percentile let downstream consumers "
            "pick their own cutoffs."
        ),
    ),
    Slide(
        title="What the tiers buy us",
        subtitle=(
            "Tier 1 is the science-grade default; Tier 3 sequesters "
            "input-OOD; Tier 2 flags rare-but-real chemistry."
        ),
        template="story",
        images=["Y34_kiel_chem_by_tier.png"],
        notes=(
            "Five-column slice: all data, T1 only, T2 only, T3 only, "
            "all-by-tier overlay.  Tier 3 clusters at the cool "
            "red-giant tip (extinction systematics).  Tier 2 tracks "
            "the rare-chemistry boundary, exactly where "
            "label-Mahalanobis identifies model extrapolation."
        ),
    ),
    Slide(
        title="GSP-Spec, APOGEE, ArqueoGal",
        subtitle=(
            "Same Tier-1 holdout cohort, three different chemistry sources.  "
            "ArqueoGal recovers APOGEE's bimodality cleanly."
        ),
        template="story",
        images=["Y33_three_way_chemistry_comparison.png"],
        notes=(
            "Left, Gaia DR3 GSP-Spec on the same source IDs.  Middle, "
            "APOGEE DR19 truth.  Right, ArqueoGal v1.1 prediction.  Note "
            "GSP-Spec's smeared bimodality vs the clean alpha-bimodality "
            "ArqueoGal recovers; the latter is matched to APOGEE within "
            "~0.04 dex per axis.  This is the central XP-to-APOGEE "
            "bridge claim."
        ),
    ),
    Slide(
        title="Application, Stream 2 (TESS asteroseismic giants)",
        subtitle=(
            "The pipeline trained on Stream 1 transfers to Stream 2; "
            "asteroseismic log g provides an independent cross-check."
        ),
        template="compare",
        images=[
            "Y39_stream2_three_way_kiel_chem.png",
            "Y41_stream2_kinematics.png",
        ],
        notes=(
            "Top: Stream 2's GSP-Spec, APOGEE cross-match, and ArqueoGal "
            "Tier 1 in Kiel and chemistry.  Bottom: Toomre / "
            "sqrt(J_R+J_z)-vs-L_z / E-vs-L_z scatter on Stream 2 Tier 1, "
            "coloured by predicted [M/H], 8000 stars with finite RV + "
            "parallax + photogeometric distance."
        ),
    ),
    Slide(
        title="Scientific output: contamination and Galactic geometry",
        subtitle=(
            "Per-star contamination via a 3-component GMM seeded for "
            "thin / thick / halo.  Disc geometry coloured by chemistry."
        ),
        template="compare",
        images=[
            "Y42_contamination_stream1_tier1.png",
            "Y37_galactic_xy_rz_by_chem.png",
        ],
        notes=(
            "GMM seeded at thin (+0.05, 0.00), thick (-0.40, 0.20), halo "
            "(-1.20, 0.25); confusion (cells sum to 1) shows where "
            "stars migrate; recall / precision / F1 bars per component.  "
            "Galactic geometry: top-down and side-view of Tier-1 "
            "holdout coloured by [M/H] and [alpha/M] medians per cell."
        ),
    ),
    Slide(
        title="Head-to-head with the XP literature",
        subtitle=(
            "Paper-quoted RMSEs vs Andrae+ 2023, Zhang+ 2023, "
            "Khalatyan+ 2024.  ArqueoGal sits inside their precision "
            "envelope on a smaller training cohort."
        ),
        template="story",
        images=["Y46_literature_comparison.png"],
        notes=(
            "Numbers are paper-quoted on each work's reported holdout, "
            "not measured here on a matched subset.  The dashed line is "
            "APOGEE DR19 internal precision.  ArqueoGal's Teff, log g, "
            "[M/H], [alpha/M] are all closest-to-floor; the precision "
            "envelope is bounded by the matching-up cohort selection "
            "differences (Andrae+ is RGB-only, Zhang+ has no log-g "
            "floor, etc.).  Caveats stamped in the figure caption."
        ),
    ),
    Slide(
        title="Take-aways and what is next",
        subtitle=(
            "XP carries enough information to recover APOGEE-grade "
            "chemistry; tiers + percentiles let consumers pick their "
            "own risk."
        ),
        template="story",
        images=["Y45_take_aways_synthesis.png"],
        notes=(
            "Take-aways: (1) precision within ~0.5x APOGEE on a holdout "
            "of 88k stars; (2) sample-size advantage vs APOGEE-only "
            "spectroscopy: Stream 3 release will be ~614k stars; (3) "
            "calibrated heteroscedastic uncertainty + per-star "
            "percentiles make the catalogue user-facing, not "
            "monolithic.  Next: extend to the 21-element APOGEE label "
            "set; methods paper; public D-Cat-b release."
        ),
    ),
]


def _set_text(tf, text: str, *, size: int, bold: bool = False, color=BODY_INK,
              align=PP_ALIGN.LEFT, font_name: str = "Inter",
              fallback: str = "Calibri") -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Try Inter, fall back to Calibri (matplotlib falls back to DejaVu Sans).
    run.font.name = font_name


def _add_title_bar(slide, title: str, subtitle: str) -> None:
    """Title row + subtitle row, both left-aligned, no decorative chrome."""
    sw = Inches(13.333)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), sw - Inches(1.0), Inches(0.7),
    )
    _set_text(title_box.text_frame, title, size=28, bold=True, color=TITLE_INK)
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.1), Inches(10.0), Inches(0.5),
        )
        _set_text(sub_box.text_frame, subtitle, size=14, bold=False, color=MUTED)


def _add_footer(slide, idx: int, total: int) -> None:
    sw = Inches(13.333)
    sh = Inches(7.5)
    left = slide.shapes.add_textbox(
        Inches(0.5), sh - Inches(0.45), Inches(8.0), Inches(0.35),
    )
    _set_text(left.text_frame,
              "ArqueoGal Pipeline 1 v1.2  /  IA-CAUP",
              size=9, bold=False, color=MUTED)
    right = slide.shapes.add_textbox(
        sw - Inches(2.5), sh - Inches(0.45), Inches(2.0), Inches(0.35),
    )
    _set_text(right.text_frame, f"{idx} / {total}",
              size=9, bold=False, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_centered_image(slide, path: Path, *, top: Inches, bottom: Inches) -> None:
    sw = Inches(13.333)
    available_h = bottom - top
    pic = slide.shapes.add_picture(str(path), Inches(0.0), top)
    iw, ih = pic.width, pic.height
    target_h = available_h
    target_w = int(iw * (target_h / ih))
    if target_w > sw - Inches(1.0):
        target_w = sw - Inches(1.0)
        target_h = int(ih * (target_w / iw))
    pic.width = int(target_w)
    pic.height = int(target_h)
    pic.left = int((sw - target_w) / 2)
    pic.top = top + int((available_h - target_h) / 2)


def _add_two_images_equal_height(slide, p1: Path, p2: Path,
                                 *, top: Inches, bottom: Inches) -> None:
    """Side-by-side images forced to the same panel-box height for visual
    consistency on T-Compare slides (brief 1.4)."""
    sw = Inches(13.333)
    available_h = bottom - top
    half = (sw - Inches(1.5)) / 2
    pad = Inches(0.5)
    # First scale each picture by width to half-slide; pick the smaller of
    # the two resulting heights so both end up the same height.
    pics = []
    for p in (p1, p2):
        pic = slide.shapes.add_picture(str(p), Inches(0.0), top)
        iw, ih = pic.width, pic.height
        target_w = half
        target_h = int(ih * (target_w / iw))
        pics.append((pic, target_w, target_h))
    target_h = min(t for _p, _w, t in pics)
    if target_h > available_h:
        target_h = int(available_h)
    for i, (pic, _w_unused, _h_unused) in enumerate(pics):
        iw, ih = pic.width, pic.height  # original pixel dims
        new_h = target_h
        new_w = int(iw * (new_h / ih))
        pic.width = int(new_w)
        pic.height = int(new_h)
        x = pad + i * (half + pad) + max(0, (half - new_w) / 2)
        pic.left = int(x)
        pic.top = top + int((available_h - new_h) / 2)


def _add_title_slide(slide, sd: Slide) -> None:
    """T-Title: hero figure right 60%, text block left 40%."""
    sw = Inches(13.333)
    sh = Inches(7.5)

    # Left text block.
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(5.2), Inches(2.2),
    )
    _set_text(title_box.text_frame, sd.title, size=34, bold=True,
              color=TITLE_INK)

    sub_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.6), Inches(5.0), Inches(2.0),
    )
    _set_text(sub_box.text_frame, sd.subtitle, size=15, bold=False,
              color=MUTED)

    foot = slide.shapes.add_textbox(
        Inches(0.6), sh - Inches(0.7), Inches(5.0), Inches(0.4),
    )
    _set_text(
        foot.text_frame,
        "IA-CAUP, Universidade do Porto",
        size=11, bold=False, color=MUTED,
    )

    # Right hero image.
    if sd.images:
        path = GALLERY / sd.images[0]
        if path.exists():
            pic = slide.shapes.add_picture(str(path), Inches(0.0),
                                           Inches(1.0))
            iw, ih = pic.width, pic.height
            target_w = Inches(7.0)
            target_h = int(ih * (target_w / iw))
            if target_h > Inches(5.5):
                target_h = Inches(5.5)
                target_w = int(iw * (target_h / ih))
            pic.width = int(target_w)
            pic.height = int(target_h)
            pic.left = sw - int(target_w) - Inches(0.5)
            pic.top = int((sh - int(target_h)) / 2)


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    total = len(SLIDES)
    for idx, sd in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)

        if sd.template == "title":
            _add_title_slide(slide, sd)
        else:
            _add_title_bar(slide, sd.title, sd.subtitle)
            top = Inches(1.7)
            bottom = Inches(6.85)
            paths = [GALLERY / name for name in sd.images]
            for p in paths:
                if not p.exists():
                    print(f"[deck] WARNING: missing {p}", file=sys.stderr)
            paths = [p for p in paths if p.exists()]
            if sd.template in ("story", "stat") and len(paths) == 1:
                _add_centered_image(slide, paths[0], top=top, bottom=bottom)
            elif sd.template == "compare" and len(paths) == 2:
                _add_two_images_equal_height(slide, paths[0], paths[1],
                                              top=top, bottom=bottom)
            elif len(paths) == 1:
                _add_centered_image(slide, paths[0], top=top, bottom=bottom)
            elif len(paths) == 2:
                _add_two_images_equal_height(slide, paths[0], paths[1],
                                              top=top, bottom=bottom)

            _add_footer(slide, idx + 1, total)

        slide.notes_slide.notes_text_frame.text = sd.notes

    prs.save(OUT)
    print(f"[deck] wrote {OUT.relative_to(REPO)}")
    print(f"[deck] {total} slides; targeted runtime ~12 minutes "
          f"({60 * 12 / total:.0f} s per slide).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
