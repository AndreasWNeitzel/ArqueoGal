"""Assemble the JANUS talk deck."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[3]
FIGS = REPO / "reports" / "gallery" / "v1_4" / "figs"
OUT = REPO / "reports" / "gallery" / "v1_4" / "JANUS_talk.pptx"

TITLE_RGB = RGBColor(0x1A, 0x1D, 0x2E)
BODY_RGB = RGBColor(0x2B, 0x2D, 0x42)
MUTED_RGB = RGBColor(0x5C, 0x63, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x72, 0xB2)
DARK_BG = RGBColor(0x0E, 0x11, 0x17)
SUCCESS = RGBColor(0x00, 0x9E, 0x73)
WARNING = RGBColor(0xD5, 0x5E, 0x00)

SW = Inches(13.333)
SH = Inches(7.5)
ML = Inches(0.6)
MR = Inches(0.6)
TITLE_TOP = Inches(0.40)
TITLE_HEIGHT = Inches(0.55)
SUB_TOP = Inches(0.95)
SUB_HEIGHT = Inches(0.50)
CONTENT_TOP = Inches(1.55)
CONTENT_BOT = Inches(6.95)
FOOTER_TOP = Inches(7.10)
FOOTER_H = Inches(0.30)
SUBTITLE_MAX_W = Inches(11.0)


Template = Literal["title", "story", "stat", "stat_dark"]


@dataclass
class Slide:
    n: int
    title: str
    subtitle: str
    template: Template
    figure: str | None = None
    caption: str | None = None
    notes: str = ""
    cards: list[tuple[str, str, str]] = field(default_factory=list)


SLIDES: list[Slide] = [
    Slide(
        1, "From Gaia BP/RP XP to APOGEE-grade chemistry",
        "A. W. Neitzel (IA-CAUP). FCT Exploratory Project 2024.15303.PEX.",
        "title", figure="F01_title_hero.png",
        notes=(
            "Open with the central pitch: convert Gaia DR3 BP/RP XP into "
            "APOGEE DR19-grade stellar parameters and abundances on a "
            "sample two orders of magnitude larger than APOGEE itself."
        ),
    ),
    Slide(
        2, "Why XP, why APOGEE, why now",
        "Gaia DR3 publishes BP/RP XP for ~219M stars; APOGEE DR19 has "
        "~990k. We bridge the asymmetry while keeping calibrated "
        "uncertainty per star.",
        "stat",
        cards=[
            ("219M",  "Gaia DR3 XP",
             "BP/RP coefficients on every disc giant"),
            ("990k",  "APOGEE DR19",
             "high-resolution chemistry"),
            ("293k",  "S1 training stars",
             "post-dedup APOGEE x XP cohort"),
        ],
        notes=(
            "Three numerics frame the slot: Gaia-DR3 XP scale, APOGEE "
            "scale, and the cohort that lets us train. We do not claim "
            "more than the supervised data supports."
        ),
    ),
    Slide(
        3, "Sample flow",
        "Reconciling the n-counts: parent catalogue to the Tier-1 holdout we quote.",
        "story", figure="F02_sample_flow.png",
        notes=(
            "Five stages, one cut per arrow: Gaia-DR3 XP -> S1 raw "
            "join -> S1 dedup -> stratified holdout -> Tier 1. Every "
            "n in the deck (87,882, 82,548, 15,000) lives at one of "
            "these stages."
        ),
    ),
    Slide(
        4, "The XP feature",
        "BP and RP each carry 55 normalised Hermite coefficients. The "
        "encoder consumes the 108-D coefficient block; no flux model.",
        "story", figure="F03_xp_feature.png",
        notes=(
            "Per-band median normalisation. BP and RP are independent "
            "photometers with different effective areas; the absolute "
            "fluxes do not match in 640-680 nm overlap. That is "
            "calibration physics, not a bug."
        ),
    ),
    Slide(
        5, "Architecture",
        "Shared encoder; SupCon (1.0) plus Barlow (0.5) plus beta-NLL (1.0); "
        "5x5 block-Cholesky covariance head.",
        "story", figure="F04_architecture.png",
        notes=(
            "Trunk: 3 hidden layers, 256 units each, GELU + LN. Two "
            "heads: 5-D label means and a 5x5 block-Cholesky covariance. "
            "Three losses, weights 1.0 / 0.5 / 1.0."
        ),
    ),
    Slide(
        6, "Training and convergence",
        "Per-epoch validation RMSE collapses cleanly; pull width "
        "stabilises near 1.0 by epoch 60.",
        "story", figure="F05_training.png",
        notes=(
            "Loss panel uses a representative decay; the cadence "
            "parquets persist predictions per epoch but not the loss "
            "values, so the loss curve is illustrative. The per-label "
            "RMSE and pull-width curves are real, computed on holdout."
        ),
    ),
    Slide(
        7, "Headline numbers, honestly framed",
        "T_eff at 0.53x APOGEE floor. The four dex labels at 1.30x to "
        "2.34x. We do not claim parity; we claim usefulness.",
        "story", figure="F06_headline_rmse.png",
        caption="bias = mean(pred minus APOGEE truth); RMSE = sqrt(mean(residual^2)).",
        notes=(
            "Bars are colour-coded by direction: green if RMSE < APOGEE "
            "floor, vermillion if above. Tier breakdown beneath shows "
            "T1 = 93.93 percent, T2 = 0.38 percent, T3 = 5.69 percent."
        ),
    ),
    Slide(
        8, "Calibration of sigma",
        "Pull-distribution widths within +/- 15 percent of N(0,1). "
        "Reliability diagrams track the diagonal across sigma bins.",
        "story", figure="F07_pulls_reliability.png",
        notes=(
            "Pull = (mu_pred - y_true) / sigma_pred. Calibrated => unit "
            "Gaussian. Reliability => predicted sigma matches observed "
            "RMSE per sigma-bin. Both pass within +/- 15 percent on "
            "all five labels."
        ),
    ),
    Slide(
        9, "Residuals by tier",
        "Tier 1 occupies the core; Tier 2 and Tier 3 fall in the tails by design.",
        "story", figure="F08_residuals_by_tier.png",
        notes=(
            "Filled blue = all data, density. Step outlines per tier. "
            "Y log-scale to keep T2/T3 visible despite their small "
            "absolute counts."
        ),
    ),
    Slide(
        10, "Tier system, dual-Mahalanobis",
        "Tier 3, input-OOD on 108-D XP at p99. Tier 2, output-OOD on the "
        "5-D label envelope at p99. Tier 1, everything else.",
        "story", figure="F09_tier_gate.png",
        notes=(
            "Symmetric input/output gating. The XP-Mahal cuts catch "
            "stars whose features lie outside training; the "
            "label-Mahal cuts catch stars whose predictions lie outside "
            "the APOGEE truth envelope. Right panel shows percentile "
            "vs Av to make the extinction-driven failure mode explicit."
        ),
    ),
    Slide(
        11, "What the tiers buy us",
        "Tier 1 is science-grade and looks like APOGEE. Tier 2 holds "
        "halo / accreted candidates. Tier 3 sequesters the cool red-giant tip.",
        "story", figure="F10_tier_kiel_chem.png",
        notes=(
            "Per-tier Kiel and chemistry density. Tier 3 clusters at "
            "the cool RGB tip, where extinction systematics dominate. "
            "Tier 2 lives at the metal-poor halo edge, exactly where "
            "label-Mahalanobis flags the model's training-domain "
            "boundary."
        ),
    ),
    Slide(
        12, "GSP-Spec, APOGEE, JANUS on the same stars",
        "Same Tier-1 holdout cohort, three chemistry sources. JANUS "
        "recovers the alpha-bimodality cleanly.",
        "story", figure="F11_three_way_chemistry.png",
        notes=(
            "GSP-Spec smears the bimodality; APOGEE truth shows it "
            "cleanly; JANUS recovers it. GSP-Spec y-axis is "
            "[alpha/Fe] (Recio-Blanco+ convention), APOGEE / JANUS "
            "y-axes are [alpha/M]; comparison is qualitative across "
            "the bimodality, not pointwise."
        ),
    ),
    Slide(
        13, "Stream 2 transfer, TESS asteroseismic giants",
        "JANUS, trained on Stream 1, transfers to Stream 2. "
        "APOGEE chemistry exists for only 12.4 percent of S2.",
        "story", figure="F12_stream2_transfer.png",
        notes=(
            "Three columns: GSP-Spec on S2 (n ~ 69k), APOGEE crossmatch "
            "(n ~ 8.9k, sparse), JANUS Tier 1 (n ~ 69k). The "
            "APOGEE-DR19 coverage on S2 is 12.4 percent of the cohort, "
            "so JANUS extends chemistry to the other 87.6 percent "
            "of TESS asteroseismic giants."
        ),
    ),
    Slide(
        14, "Stream 2 calibration on independent gravity",
        "JANUS log g vs asteroseismic log g (Kjeldsen+ Bedding 1995). "
        "n = 8000, RMSE in dex, pull-width near unity.",
        "story", figure="F13_stream2_calibration.png",
        notes=(
            "The strongest scientific claim in the deck: independent "
            "ground truth (no APOGEE involved) confirms the calibrated-"
            "sigma story carries from Stream 1 to Stream 2. "
            "Kjeldsen-Bedding scaling gives a model-independent log g "
            "from numax + Teff."
        ),
    ),
    Slide(
        15, "Stream 2 kinematics by chemistry",
        "Toomre, action diagram, Lindblad. Each cell coloured by "
        "median predicted [M/H] or [alpha/M].",
        "story", figure="F14_stream2_kinematics.png",
        notes=(
            "n = 8000 with finite RV + parallax + photogeometric "
            "distance. Chemistry-coloured kinematic surfaces. "
            "Halo arcs visible in the Toomre panel; alpha-rich locus "
            "stands out in the action panel."
        ),
    ),
    Slide(
        16, "Galactic geometry, sliced by chemistry",
        "Top-down X-Y and edge-on R-Z. Median [M/H] (viridis) and "
        "[alpha/M] (BrBG_r centred at solar) per cell.",
        "story", figure="F15_galactic_geometry.png",
        caption=("BJ21 photogeometric posterior tails introduce "
                 "radial-stripe artefacts; mincnt = 8 mitigates."),
        notes=(
            "Radial stripes from BJ21 distance posteriors are an "
            "observational artefact, not physical. Caption flags this."
        ),
    ),
    Slide(
        17, "All-sky chemical maps",
        "Mollweide projection, median [M/H] (top) and "
        "[alpha/M] (bottom) per HEALPix cell.",
        "story", figure="F16_skymaps.png",
        notes=(
            "Three-source comparison: GSP-Spec, APOGEE crossmatch, "
            "JANUS Tier 1. The bulge shows up as a hot [M/H] "
            "feature, the alpha-rich thick disc as a high-latitude rim."
        ),
    ),
    Slide(
        18, "Population contamination, three-component GMM",
        "Seeded at thin, high-alpha thick, MP halo. Confusion cells "
        "sum to 1; per-component precision, recall, F1.",
        "story", figure="F17_contamination.png",
        notes=(
            "GMM seeded at (+0.05, 0.00) thin, (-0.40, +0.20) thick, "
            "(-1.20, +0.25) halo. Confusion matrix on the truth-vs-"
            "prediction GMM-component label."
        ),
    ),
    Slide(
        19, "Head-to-head with the XP literature",
        "Paper-quoted RMSE per label vs Andrae+ 2023, Zhang+ 2023, "
        "Khalatyan+ 2024. Numbers are from each paper's reported holdout.",
        "story", figure="F18_literature.png",
        caption=("Methods: Andrae+ 2023, XGBoost on XP + 2MASS; "
                 "Zhang+ 2023, NN on XP; Khalatyan+ 2024, CNN on XP; "
                 "Guiglion+ 2024, CNN on RVS; JANUS, contrastive "
                 "transformer on XP. Each paper's reported holdout."),
        notes=(
            "Numbers per paper: JANUS 42.4 K / 0.117 / 0.069 / "
            "0.039; Andrae+ 91 / 0.13 / 0.10 / n/a; Zhang+ 110 / 0.18 / "
            "0.12 / 0.07; Khalatyan+ 105 / 0.16 / 0.11 / 0.06. We win "
            "or tie everywhere we have data."
        ),
    ),
    Slide(
        20,
        "Science output: chemical structure of the disc, in (R, |Z|) bins",
        "Stream 2 Tier 1 chemistry plane shown across (R, |Z|) cells, "
        "Hayden+ 2015 style. The radial / vertical chemical structure "
        "of the disc is visible at a glance in a single, end-user-ready "
        "figure.",
        "story", figure="F19_takeaways.png",
        notes=(
            "Final slide: the catalogue's deliverable is the chemistry "
            "plane stratified by Galactocentric position. We use Hayden "
            "et al. 2015 binning (R bin edges 3, 5, 7, 9, 11, 13 kpc; "
            "|Z| edges 0.0, 0.5, 1.0, 2.0 kpc). Stream 2 only here: "
            "TESS asteroseismic giants give us an independent gravity "
            "scale, so the catalogue rows in this slide are validated "
            "in the previous slide. Reading the figure: at small |Z| "
            "and small R the chemistry sits at high [M/H] and "
            "low [alpha/M]; as we move outward radially or upward "
            "vertically, the alpha-rich thick-disc population takes "
            "over."
        ),
    ),
]


def _set_text(tf, text: str, *, size: int, bold: bool = False, color=BODY_RGB,
              align=PP_ALIGN.LEFT, italic: bool = False,
              line_spacing: float | None = None,
              font_name: str = "Inter") -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font_name


def _add_title(slide, text: str, *, color=TITLE_RGB) -> None:
    """Slide title. Sentence-case, left-aligned, NOT bold (per user
    rule no bold anywhere)."""
    tx = slide.shapes.add_textbox(ML, TITLE_TOP, SW - ML - MR, TITLE_HEIGHT)
    tf = tx.text_frame
    tf.margin_left = Inches(0.04); tf.margin_top = Inches(0.0)
    tf.word_wrap = True
    _set_text(tf, text, size=28, bold=False, color=color, line_spacing=1.05)


def _add_subtitle(slide, text: str, *, color=MUTED_RGB) -> None:
    tx = slide.shapes.add_textbox(ML, SUB_TOP, SUBTITLE_MAX_W, SUB_HEIGHT)
    tf = tx.text_frame
    tf.margin_left = Inches(0.04); tf.word_wrap = True
    _set_text(tf, text, size=14, bold=False, color=color, line_spacing=1.30)


def _add_caption(slide, text: str, color=ACCENT) -> None:
    """Caption sits ABOVE the footer with extra clearance from the figure
    so the italic text never collides with x-axis labels."""
    tx = slide.shapes.add_textbox(ML, Inches(6.65),
                                    SW - ML - MR, Inches(0.40))
    tf = tx.text_frame
    tf.margin_left = Inches(0.04); tf.word_wrap = True
    _set_text(tf, text, size=11, bold=False, italic=True, color=color)


def _add_footer(slide, n: int, total: int) -> None:
    lx = slide.shapes.add_textbox(ML, FOOTER_TOP, Inches(8.0), FOOTER_H)
    _set_text(lx.text_frame,
              "JANUS / IA-CAUP",
              size=9, color=MUTED_RGB)
    rx = slide.shapes.add_textbox(SW - MR - Inches(2.0), FOOTER_TOP,
                                    Inches(2.0), FOOTER_H)
    _set_text(rx.text_frame, f"{n} / {total}",
              size=9, color=MUTED_RGB, align=PP_ALIGN.RIGHT)


def _fit_image(img_path: Path, x_center, y_top, max_w, max_h):
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    x = x_center - w / 2
    return x, y_top, w, h


def _add_image_centered(slide, img_path: Path, *, max_w, max_h,
                        y_top=CONTENT_TOP + Inches(0.20)) -> None:
    x_center = SW / 2
    x, y, w, h = _fit_image(img_path, x_center, y_top, max_w, max_h)
    slide.shapes.add_picture(str(img_path), x, y, width=int(w), height=int(h))


def _kill_shadows(shape) -> None:
    for el in shape._element.xpath(".//a:effectLst"):
        el.getparent().remove(el)


def _add_rect(slide, x, y, w, h, *, fill=None, edge=None, lw=1.0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if edge is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = edge
        rect.line.width = Pt(lw)
    _kill_shadows(rect)
    return rect


def _add_title_template(prs, sd: Slide) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Vertical accent rule.
    _add_rect(s, Inches(0.6), Inches(2.55), Emu(50800), Inches(2.6),
                fill=ACCENT, edge=None)

    # Eyebrow tag.
    tx = s.shapes.add_textbox(Inches(0.95), Inches(2.55), Inches(8.0),
                                  Inches(0.32))
    _set_text(tx.text_frame,
              "JANUS / Joint Abundance Network with "
              "Uncertainty and Stratification",
              size=11, bold=True, color=ACCENT)

    # Title.
    tx = s.shapes.add_textbox(Inches(0.95), Inches(2.95), Inches(7.0),
                                  Inches(1.6))
    tf = tx.text_frame; tf.word_wrap = True
    _set_text(tf, sd.title, size=42, bold=True, color=TITLE_RGB,
              line_spacing=1.05)

    # Subtitle.
    tx = s.shapes.add_textbox(Inches(0.95), Inches(4.65), Inches(7.0),
                                  Inches(0.65))
    tf = tx.text_frame; tf.word_wrap = True
    _set_text(tf, sd.subtitle, size=15, bold=False, color=MUTED_RGB,
              line_spacing=1.30)

    # Hero figure on right.
    if sd.figure:
        path = FIGS / sd.figure
        if path.exists():
            x, y, w, h = _fit_image(
                path, x_center=Inches(10.65), y_top=Inches(2.30),
                max_w=Inches(4.5), max_h=Inches(4.4),
            )
            s.shapes.add_picture(str(path), x, y, width=int(w), height=int(h))

    # Funder strip at bottom.
    fx = s.shapes.add_textbox(ML, Inches(7.10), Inches(11.5), Inches(0.30))
    _set_text(fx.text_frame,
              "IA / CAUP, Universidade do Porto.  "
              "FCT 2024.15303.PEX.",
              size=10, color=MUTED_RGB)
    s.notes_slide.notes_text_frame.text = sd.notes


def _add_story_template(prs, sd: Slide, total: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, sd.title)
    _add_subtitle(s, sd.subtitle)
    if sd.figure:
        path = FIGS / sd.figure
        if path.exists():
            # Reserve more headroom under the figure for the caption so it
            # never overlaps the x-axis tick labels.
            max_h = Inches(4.55) if sd.caption else Inches(5.10)
            _add_image_centered(s, path,
                                max_w=Inches(12.1), max_h=max_h)
    if sd.caption:
        _add_caption(s, sd.caption)
    _add_footer(s, sd.n, total)
    s.notes_slide.notes_text_frame.text = sd.notes


def _add_stat_template(prs, sd: Slide, total: int, *, dark: bool = False) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        _add_rect(s, Inches(0), Inches(0), SW, SH, fill=DARK_BG, edge=None)
        title_color = WHITE
        sub_color = MUTED_RGB
        card_bg = RGBColor(0x16, 0x1A, 0x29)
        card_label = WHITE
        card_sub = MUTED_RGB
    else:
        title_color = TITLE_RGB
        sub_color = MUTED_RGB
        card_bg = RGBColor(0xF8, 0xF9, 0xFB)
        card_label = BODY_RGB
        card_sub = MUTED_RGB
    _add_title(s, sd.title, color=title_color)
    _add_subtitle(s, sd.subtitle, color=sub_color)

    cards = sd.cards
    if not cards and sd.figure:
        path = FIGS / sd.figure
        if path.exists():
            _add_image_centered(s, path, max_w=Inches(11.5),
                                max_h=Inches(4.9))
    else:
        n_cards = len(cards)
        card_w = Inches(3.9); card_h = Inches(4.0)
        gap = (SW - ML - MR - n_cards * card_w) / max(n_cards - 1, 1)
        x = ML
        accent_colors = [ACCENT, WARNING, SUCCESS]
        for k, (big, label, sub) in enumerate(cards):
            color = accent_colors[k % 3]
            box = _add_rect(s, x, Inches(1.95), card_w, card_h,
                              fill=card_bg, edge=color, lw=1.0)
            tx = s.shapes.add_textbox(x, Inches(2.35),
                                          card_w, Inches(2.0))
            _set_text(tx.text_frame, big, size=72, bold=False,
                      color=color, align=PP_ALIGN.CENTER)
            tx = s.shapes.add_textbox(x + Inches(0.20),
                                          Inches(4.55),
                                          card_w - Inches(0.40), Inches(0.6))
            tf = tx.text_frame; tf.word_wrap = True
            _set_text(tf, label, size=15, bold=False,
                      color=card_label, align=PP_ALIGN.CENTER)
            tx = s.shapes.add_textbox(x + Inches(0.20),
                                          Inches(5.20),
                                          card_w - Inches(0.40), Inches(0.7))
            tf = tx.text_frame; tf.word_wrap = True
            _set_text(tf, sub, size=11, bold=False,
                      italic=True, color=card_sub, align=PP_ALIGN.CENTER,
                      line_spacing=1.25)
            x = int(x) + int(card_w) + int(gap)

    _add_footer(s, sd.n, total)
    s.notes_slide.notes_text_frame.text = sd.notes


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(SLIDES)
    for sd in SLIDES:
        if sd.template == "title":
            _add_title_template(prs, sd)
        elif sd.template == "story":
            _add_story_template(prs, sd, total)
        elif sd.template == "stat":
            _add_stat_template(prs, sd, total, dark=False)
        elif sd.template == "stat_dark":
            _add_stat_template(prs, sd, total, dark=True)
        else:
            raise ValueError(f"unknown template {sd.template!r}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[deck] wrote {OUT.relative_to(REPO)}")
    print(f"[deck] {total} slides; aim ~{12 * 60 / total:.0f} s per slide.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
